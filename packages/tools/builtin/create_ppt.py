"""create_ppt — generate a PowerPoint presentation from structured slides."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

from pydantic import BaseModel, Field

from citnega.packages.protocol.callables.base import BaseCallable
from citnega.packages.protocol.callables.types import CallableType
from citnega.packages.tools.builtin._tool_base import ToolOutput, tool_policy

if TYPE_CHECKING:
    from citnega.packages.protocol.callables.context import CallContext


class Slide(BaseModel):
    title: str = Field(description="Slide title.")
    bullets: list[str] = Field(default_factory=list, description="Bullet points for the slide body.")
    notes: str = Field(default="", description="Speaker notes for this slide.")


class CreatePPTInput(BaseModel):
    title: str = Field(description="Presentation title (shown on the first cover slide).")
    slides: list[Slide] = Field(description="Ordered list of slides.")
    filename: str = Field(description="Output file path, e.g. ~/Desktop/deck.pptx")
    subtitle: str = Field(default="", description="Subtitle text on the cover slide.")


class CreatePPTTool(BaseCallable):
    """Generate a PowerPoint (.pptx) presentation from a title and list of slides."""

    name = "create_ppt"
    description = (
        "Create a PowerPoint (.pptx) presentation from a title and list of slides. "
        "Each slide has a title, bullet points, and optional speaker notes. "
        "A cover slide is added automatically. Returns the path of the created file."
    )
    callable_type = CallableType.TOOL
    input_schema = CreatePPTInput
    output_schema = ToolOutput
    policy = tool_policy(
        timeout_seconds=30.0,
        requires_approval=False,
        network_allowed=False,
    )

    async def _execute(self, input: CreatePPTInput, context: CallContext) -> ToolOutput:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.util import Inches, Pt  # type: ignore[import-untyped]
        except ImportError:
            return ToolOutput(result="[create_ppt: python-pptx not installed — run: pip install python-pptx]")

        out_path = Path(input.filename).expanduser().resolve()
        out_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()

        # Cover slide (layout 0 = title slide)
        cover_layout = prs.slide_layouts[0]
        cover = prs.slides.add_slide(cover_layout)
        cover.shapes.title.text = input.title
        if input.subtitle and cover.placeholders[1]:
            cover.placeholders[1].text = input.subtitle

        # Content slides (layout 1 = title + content)
        content_layout = prs.slide_layouts[1]
        for slide_data in input.slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data.title

            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(slide_data.bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                    tf.paragraphs[0].level = 0
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

            if slide_data.notes:
                slide.notes_slide.notes_text_frame.text = slide_data.notes

        try:
            prs.save(str(out_path))
        except Exception as exc:
            return ToolOutput(result=f"[create_ppt: failed to write: {exc}]")

        return ToolOutput(
            result=f"PPTX created: {out_path}  ({len(input.slides) + 1} slides, {out_path.stat().st_size} bytes)"
        )
